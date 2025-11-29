"""
Document parsing module for extracting text and figures from various formats
"""

import os
import re
import zipfile
import tempfile
from typing import Dict, List, Optional, Tuple
from pathlib import Path

# PDF parsing
from PyPDF2 import PdfReader

# DOCX parsing
from docx import Document as DocxDocument
from docx.shared import Inches

# PPTX parsing
from pptx import Presentation
from pptx.util import Inches as PptxInches

# LaTeX parsing
from pylatexenc.latex2text import LatexNodes2Text

class DocumentParser:
    """Base class for document parsing"""
    
    def __init__(self, file_path: str, output_dir: str):
        self.file_path = file_path
        self.output_dir = output_dir
        self.figures_dir = os.path.join(output_dir, "figures")
        os.makedirs(self.figures_dir, exist_ok=True)
    
    def parse(self) -> Dict:
        """Parse the document and return extracted content"""
        raise NotImplementedError


class PDFParser(DocumentParser):
    """Parse PDF documents"""
    
    def _clean_text(self, text: str) -> str:
        """Clean text to ensure valid UTF-8 encoding"""
        if not text:
            return ""
        # Remove null bytes and other problematic characters
        text = text.replace('\x00', '')
        # Encode and decode to handle invalid UTF-8
        text = text.encode('utf-8', errors='ignore').decode('utf-8', errors='ignore')
        # Remove other control characters except newlines and tabs
        cleaned = ''.join(char for char in text if char.isprintable() or char in '\n\t\r')
        return cleaned
    
    def parse(self) -> Dict:
        text_content = []
        sections = []
        figures = []
        
        try:
            reader = PdfReader(self.file_path)
            
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                page_text = self._clean_text(page_text)
                text_content.append(page_text)
                
                # Try to extract images (safely)
                try:
                    if hasattr(page, 'images'):
                        for img_num, image in enumerate(page.images):
                            try:
                                img_path = os.path.join(
                                    self.figures_dir,
                                    f"figure_p{page_num}_{img_num}.{image.name.split('.')[-1] if '.' in image.name else 'png'}"
                                )
                                with open(img_path, 'wb') as f:
                                    f.write(image.data)
                                figures.append(img_path)
                            except Exception:
                                pass
                except Exception:
                    # Skip image extraction if it fails
                    pass
            
            # Extract sections based on common patterns
            full_text = "\n".join(text_content)
            sections = self._extract_sections(full_text)
            
            return {
                "text": full_text,
                "sections": sections,
                "figures": figures,
                "page_count": len(reader.pages)
            }
            
        except Exception as e:
            return {
                "text": "",
                "sections": [],
                "figures": [],
                "error": str(e)
            }
    
    def _extract_sections(self, text: str) -> List[Dict]:
        """Extract sections from text based on common academic paper patterns"""
        sections = []
        
        # Common section headers
        section_patterns = [
            r'^(?:Abstract|ABSTRACT)',
            r'^(?:1\.?\s*)?(?:Introduction|INTRODUCTION)',
            r'^(?:2\.?\s*)?(?:Related Work|RELATED WORK|Background|BACKGROUND|Literature Review)',
            r'^(?:3\.?\s*)?(?:Method(?:ology)?|METHODOLOGY?|Approach|APPROACH)',
            r'^(?:4\.?\s*)?(?:Experiment(?:s)?|EXPERIMENTS?|Results|RESULTS)',
            r'^(?:5\.?\s*)?(?:Discussion|DISCUSSION)',
            r'^(?:6\.?\s*)?(?:Conclusion(?:s)?|CONCLUSIONS?)',
            r'^(?:References|REFERENCES|Bibliography)'
        ]
        
        # Try to find sections
        lines = text.split('\n')
        current_section = {"title": "Introduction", "content": []}
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            is_section_header = False
            for pattern in section_patterns:
                if re.match(pattern, line, re.IGNORECASE):
                    if current_section["content"]:
                        sections.append({
                            "title": current_section["title"],
                            "content": "\n".join(current_section["content"])
                        })
                    current_section = {"title": line, "content": []}
                    is_section_header = True
                    break
            
            if not is_section_header:
                current_section["content"].append(line)
        
        if current_section["content"]:
            sections.append({
                "title": current_section["title"],
                "content": "\n".join(current_section["content"])
            })
        
        return sections


class DOCXParser(DocumentParser):
    """Parse DOCX documents"""
    
    def parse(self) -> Dict:
        text_content = []
        sections = []
        figures = []
        
        try:
            doc = DocxDocument(self.file_path)
            
            # Extract paragraphs
            for para in doc.paragraphs:
                text_content.append(para.text)
            
            # Extract images
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        img_data = rel.target_part.blob
                        img_ext = rel.target_ref.split('.')[-1]
                        img_path = os.path.join(
                            self.figures_dir,
                            f"figure_{len(figures)}.{img_ext}"
                        )
                        with open(img_path, 'wb') as f:
                            f.write(img_data)
                        figures.append(img_path)
                    except Exception:
                        pass
            
            # Extract sections based on headings
            current_section = {"title": "Document Start", "content": []}
            
            for para in doc.paragraphs:
                if para.style and 'Heading' in para.style.name:
                    if current_section["content"]:
                        sections.append({
                            "title": current_section["title"],
                            "content": "\n".join(current_section["content"])
                        })
                    current_section = {"title": para.text, "content": []}
                else:
                    current_section["content"].append(para.text)
            
            if current_section["content"]:
                sections.append({
                    "title": current_section["title"],
                    "content": "\n".join(current_section["content"])
                })
            
            return {
                "text": "\n".join(text_content),
                "sections": sections,
                "figures": figures
            }
            
        except Exception as e:
            return {
                "text": "",
                "sections": [],
                "figures": [],
                "error": str(e)
            }


class PPTXParser(DocumentParser):
    """Parse PPTX documents"""
    
    def parse(self) -> Dict:
        text_content = []
        sections = []
        figures = []
        
        try:
            prs = Presentation(self.file_path)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                slide_title = f"Slide {slide_num + 1}"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    slide_text.append(run.text)
                        
                        # Check if it's a title shape
                        if shape.shape_type and "TITLE" in str(shape.shape_type):
                            slide_title = shape.text or slide_title
                    
                    # Extract images
                    if shape.shape_type and "PICTURE" in str(shape.shape_type):
                        try:
                            image = shape.image
                            img_path = os.path.join(
                                self.figures_dir,
                                f"slide{slide_num}_img{len(figures)}.{image.ext}"
                            )
                            with open(img_path, 'wb') as f:
                                f.write(image.blob)
                            figures.append(img_path)
                        except Exception:
                            pass
                
                slide_content = " ".join(slide_text)
                text_content.append(slide_content)
                sections.append({
                    "title": slide_title,
                    "content": slide_content
                })
            
            return {
                "text": "\n\n".join(text_content),
                "sections": sections,
                "figures": figures,
                "slide_count": len(prs.slides)
            }
            
        except Exception as e:
            return {
                "text": "",
                "sections": [],
                "figures": [],
                "error": str(e)
            }


class LaTeXParser(DocumentParser):
    """Parse LaTeX documents (from ZIP archives)"""
    
    def parse(self) -> Dict:
        text_content = []
        sections = []
        figures = []
        
        try:
            # Extract ZIP file
            with tempfile.TemporaryDirectory() as temp_dir:
                with zipfile.ZipFile(self.file_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
                
                # Find .tex files
                tex_files = list(Path(temp_dir).rglob("*.tex"))
                
                # Find the main tex file (usually contains \documentclass)
                main_tex = None
                for tex_file in tex_files:
                    with open(tex_file, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                        if '\\documentclass' in content:
                            main_tex = tex_file
                            break
                
                if not main_tex and tex_files:
                    main_tex = tex_files[0]
                
                if main_tex:
                    with open(main_tex, 'r', encoding='utf-8', errors='ignore') as f:
                        latex_content = f.read()
                    
                    # Convert LaTeX to plain text
                    converter = LatexNodes2Text()
                    plain_text = converter.latex_to_text(latex_content)
                    text_content.append(plain_text)
                    
                    # Extract sections from LaTeX
                    sections = self._extract_latex_sections(latex_content)
                    
                    # Copy image files
                    image_extensions = ['.png', '.jpg', '.jpeg', '.pdf', '.eps']
                    for ext in image_extensions:
                        for img_file in Path(temp_dir).rglob(f"*{ext}"):
                            try:
                                img_path = os.path.join(
                                    self.figures_dir,
                                    img_file.name
                                )
                                with open(img_file, 'rb') as src:
                                    with open(img_path, 'wb') as dst:
                                        dst.write(src.read())
                                figures.append(img_path)
                            except Exception:
                                pass
            
            return {
                "text": "\n".join(text_content),
                "sections": sections,
                "figures": figures
            }
            
        except Exception as e:
            return {
                "text": "",
                "sections": [],
                "figures": [],
                "error": str(e)
            }
    
    def _extract_latex_sections(self, content: str) -> List[Dict]:
        """Extract sections from LaTeX content"""
        sections = []
        converter = LatexNodes2Text()
        
        # Pattern to match section commands
        section_pattern = r'\\(?:section|subsection|chapter)\*?\{([^}]+)\}'
        
        # Find all section matches with their positions
        matches = list(re.finditer(section_pattern, content))
        
        for i, match in enumerate(matches):
            section_title = match.group(1)
            start_pos = match.end()
            
            # Find end position (next section or end of document)
            if i + 1 < len(matches):
                end_pos = matches[i + 1].start()
            else:
                end_pos = len(content)
            
            section_content = content[start_pos:end_pos]
            
            # Convert section content to plain text
            try:
                plain_content = converter.latex_to_text(section_content)
            except:
                plain_content = section_content
            
            sections.append({
                "title": section_title,
                "content": plain_content.strip()
            })
        
        return sections


def get_parser(file_path: str, file_type: str, output_dir: str) -> DocumentParser:
    """Factory function to get appropriate parser"""
    parsers = {
        'pdf': PDFParser,
        'docx': DOCXParser,
        'pptx': PPTXParser,
        'latex': LaTeXParser,
        'zip': LaTeXParser
    }
    
    parser_class = parsers.get(file_type.lower())
    if not parser_class:
        raise ValueError(f"Unsupported file type: {file_type}")
    
    return parser_class(file_path, output_dir)


def parse_document(file_path: str, file_type: str, output_dir: str) -> Dict:
    """Main function to parse a document"""
    parser = get_parser(file_path, file_type, output_dir)
    return parser.parse()
